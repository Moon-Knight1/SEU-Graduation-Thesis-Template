# -*- coding: utf-8 -*-
"""
答辩PPT生成脚本
基于东南大学模板创建答辩PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import os

# 模板路径
TEMPLATE_PATH = r"E:\SEU-Graduation-Thesis-Template\doc\SEU01东南大学—东南集市.pptx"
OUTPUT_PATH = r"E:\SEU-Graduation-Thesis-Template\output\答辩PPT.pptx"

def load_template():
    """加载PPT模板"""
    prs = Presentation(TEMPLATE_PATH)
    return prs

def save_ppt(prs, path):
    """保存PPT文件"""
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"PPT已保存到: {path}")

def clear_all_slides(prs):
    """清除模板中的所有幻灯片"""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def create_cover_slide(prs, src_shapes_xml=None):
    """创建封面页

    模板的布局0不含形状，因此通过复制模板第0页（已有全部封面形状）来创建封面。
    形状索引基于模板分析结果：
      1  - 中文标题
      8  - 英文标题（AUTO_SHAPE）
      11 - 答辩学生
      13 - 专业班级
      15 - 指导老师

    Args:
        prs: Presentation对象
        src_shapes_xml: 预先复制的形状XML列表，如果为None则从模板第0页获取
    """
    if src_shapes_xml is None:
        src_shapes_xml = [deepcopy(shape._element) for shape in prs.slides[0].shapes]

    # 每次创建幻灯片时深拷贝形状XML，避免lxml元素父子关系冲突
    shapes_copy = [deepcopy(xml) for xml in src_shapes_xml]

    # 使用空白布局创建新幻灯片，再将封面形状复制进来
    slide_layout = prs.slide_layouts[7]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    for shape_xml in shapes_copy:
        slide.shapes._spTree.append(shape_xml)

    # 设置论文标题
    title = slide.shapes[1]  # 标题文本框
    title.text_frame.text = "基于B/S架构的企业数字化运营平台设计与实现"

    # 设置英文标题
    subtitle = slide.shapes[8]  # 英文标题文本框
    subtitle.text_frame.text = "Design and Implementation of Enterprise Digital Operation Platform Based on B/S Architecture"

    # 设置答辩学生
    student = slide.shapes[11]
    student.text_frame.text = "答辩学生：XXX"

    # 设置专业班级
    major = slide.shapes[13]
    major.text_frame.text = "专业班级：软件工程XX班"

    # 设置指导老师
    teacher = slide.shapes[15]
    teacher.text_frame.text = "指导老师：XXX"

    return slide

def create_toc_slide(prs, src_shapes_xml=None):
    """创建目录页

    通过复制模板第1页（目录页）来创建目录。
    模板目录页共28个形状，5个章节按垂直位置从上到下排列：
      18,19,20 - 第1章：编号、中文标题、英文标题
       5, 6, 7 - 第2章
       8, 9,10 - 第3章
      11,12,13 - 第4章
      14,15,16 - 第5章
    其他形状：0=背景组, 1=装饰, 2="目录", 3="CONTENTS",
              4=装饰, 17/21-25=分隔线, 26=页脚, 27=图片

    Args:
        prs: Presentation对象
        src_shapes_xml: 预先复制的形状XML列表，如果为None则从模板第1页获取
    """
    # 目录章节数据：(中文标题, 英文标题)
    toc_sections = [
        ("项目概述与系统架构", "Project Overview and Architecture"),
        ("核心功能模块设计", "Core Function Modules"),
        ("技术亮点与创新", "Technical Highlights"),
        ("系统测试与性能优化", "Testing and Optimization"),
        ("总结与展望", "Conclusion and Prospects"),
    ]

    # 各章节在模板形状中的索引：(编号shape, 中文标题shape, 英文标题shape)
    section_indices = [
        (18, 19, 20),  # 第1章 - 最上方
        (5,  6,  7),   # 第2章
        (8,  9,  10),  # 第3章
        (11, 12, 13),  # 第4章
        (14, 15, 16),  # 第5章 - 最下方
    ]

    if src_shapes_xml is None:
        src_shapes_xml = [deepcopy(shape._element) for shape in prs.slides[1].shapes]

    # 每次创建幻灯片时深拷贝形状XML，避免lxml元素父子关系冲突
    shapes_copy = [deepcopy(xml) for xml in src_shapes_xml]

    # 使用空白布局创建目录页，再将形状复制进来
    slide_layout = prs.slide_layouts[7]
    slide = prs.slides.add_slide(slide_layout)
    for shape_xml in shapes_copy:
        slide.shapes._spTree.append(shape_xml)

    # 更新各章节文本
    for i, (cn_title, en_title) in enumerate(toc_sections):
        num_idx, cn_idx, en_idx = section_indices[i]

        # 设置编号
        num_shape = slide.shapes[num_idx]
        num_shape.text_frame.text = f"0{i + 1}"

        # 设置中文标题
        cn_shape = slide.shapes[cn_idx]
        cn_shape.text_frame.text = cn_title

        # 设置英文标题
        en_shape = slide.shapes[en_idx]
        en_shape.text_frame.text = en_title

    return slide

def create_content_slide(prs, title_cn, title_en, section_num,
                         content_items=None, src_shapes_xml=None):
    """创建内容页（章节标题页）

    通过复制模板第2页（第1章章节标题页）来创建内容页。
    模板章节标题页共11个形状，结构一致：
      0  - 中文标题文本框
      1  - 英文标题文本框
      2  - 章节编号文本框（"01"、"02"等大字）
      3  - 底部装饰线
      4  - 底部装饰矩形
      5  - 顶部装饰线
      6  - 顶部装饰矩形
      7  - 页脚组合
      8  - 页眉文字（"止于至善"）
      9  - 校徽图片
      10 - 右上角装饰组合

    Args:
        prs: Presentation对象
        title_cn: 中文标题字符串
        title_en: 英文标题字符串
        section_num: 章节编号字符串（如 "01"、"02"）
        content_items: 内容条目列表（可选），若提供则创建包含内容文本的文本框
        src_shapes_xml: 预先复制的形状XML列表，如果为None则从模板第2页获取

    Returns:
        创建的幻灯片对象
    """
    if src_shapes_xml is None:
        src_shapes_xml = [deepcopy(shape._element) for shape in prs.slides[2].shapes]

    # 每次创建幻灯片时深拷贝形状XML，避免lxml元素父子关系冲突
    shapes_copy = [deepcopy(xml) for xml in src_shapes_xml]

    # 使用空白布局创建新幻灯片，再将形状复制进来
    slide_layout = prs.slide_layouts[7]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    for shape_xml in shapes_copy:
        slide.shapes._spTree.append(shape_xml)

    # 设置中文标题
    slide.shapes[0].text_frame.text = title_cn

    # 设置英文标题
    slide.shapes[1].text_frame.text = title_en

    # 设置章节编号
    slide.shapes[2].text_frame.text = section_num

    # 如果提供了内容条目，在页面下方创建内容文本框
    if content_items:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(3.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_after = Pt(6)

    return slide


def create_project_overview_slide(prs, src_shapes_xml=None):
    """创建项目概述页"""
    content_items = [
        "背景：汽车电子软件研发企业，项目过程数据分散，统计分析依赖人工",
        "目标：构建数字化运营平台，整合项目、组织、人员、工时、质量数据",
        "用户角色：部门经理、项目经理、项目成员、质量人员",
        "业务范围：项目管理、公司运营管理、系统管理"
    ]

    return create_content_slide(
        prs,
        title_cn="项目概述",
        title_en="Project Overview",
        section_num="01",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_architecture_slide(prs, src_shapes_xml=None):
    """创建系统架构页"""
    content_items = [
        "五层架构：前端展示层、请求处理层、业务逻辑层、数据访问层、数据层",
        "技术栈：Spring Boot + Vue + MySQL + Redis + MyBatis-Plus",
        "前后端分离：B/S架构，RESTful API",
        "安全框架：Shiro + JWT + LDAP统一认证"
    ]

    return create_content_slide(
        prs,
        title_cn="系统架构设计",
        title_en="System Architecture Design",
        section_num="01",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_bu_hour_function_slide(prs, src_shapes_xml=None):
    """创建BU工时功能页"""
    content_items = [
        "需求：部门经理查看各BU工时投入分布，支持按年月、组织筛选",
        "核心功能：饼图展示BU工时占比、BU详情（总工时、平均工时、趋势图）",
        "设计：WorkHourController → WorkHourService → WorkHourMapper",
        "数据流转：前端查询 → 参数校验 → 组织树遍历 → 工时聚合 → 饼图渲染"
    ]

    return create_content_slide(
        prs,
        title_cn="BU工时分布模块",
        title_en="BU Work Hour Distribution",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_bu_hour_optimization_slide(prs, src_shapes_xml=None):
    """创建BU工时优化页"""
    content_items = [
        "问题：DATE_FORMAT函数导致work_date索引失效，全表扫描",
        "方案：将月份筛选改写为日期范围查询，避免在索引字段上使用函数",
        "效果：访问类型从ALL变为range，成功命中idx_wh_work_date索引",
        "缓存策略：Redis缓存（当月1小时，历史24小时，随机抖动防雪崩）"
    ]

    return create_content_slide(
        prs,
        title_cn="BU工时模块优化",
        title_en="BU Work Hour Optimization",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_project_info_slide(prs, src_shapes_xml=None):
    """创建项目信息管理页"""
    content_items = [
        "需求：项目经理维护项目基础信息，记录交付历史",
        "核心功能：项目CRUD、条件筛选、Excel导出、交付历史追溯",
        "设计：主表+历史表双写策略，保证数据可追溯",
        "权限控制：@RequiresPermissions注解，Shiro框架"
    ]

    return create_content_slide(
        prs,
        title_cn="项目基本信息管理",
        title_en="Project Information Management",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_bug_management_function_slide(prs, src_shapes_xml=None):
    """创建缺陷管理功能页"""
    content_items = [
        "需求：质量人员统一查看Trinity、Jira、BugClose三个平台的缺陷数据",
        "核心功能：缺陷分页查询、项目聚类统计、健康度分析、飞书预警",
        "设计：BugController → BugService → BugMapper",
        "数据源：三张同步表分别存储不同平台的缺陷数据"
    ]

    return create_content_slide(
        prs,
        title_cn="缺陷管理模块",
        title_en="Bug Management Module",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_bug_management_optimization_slide(prs, src_shapes_xml=None):
    """创建缺陷管理优化页"""
    content_items = [
        "问题1：if-else分支处理三个数据源，违反开闭原则",
        "方案1：策略模式+模板方法模式重构，定义BugPageQueryStrategy接口",
        "问题2：深分页查询LIMIT offset, N效率低下",
        "方案2：延迟关联优化，先用覆盖索引查询主键，再关联获取完整数据",
        "效果：5K页从35736ms降至73.2ms，提升488倍"
    ]

    return create_content_slide(
        prs,
        title_cn="缺陷管理优化",
        title_en="Bug Management Optimization",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_serious_problem_function_slide(prs, src_shapes_xml=None):
    """创建重大问题管理功能页"""
    content_items = [
        "需求：跟踪项目执行过程中的重大问题，支持全生命周期管理",
        "核心功能：问题登记、处理、上升、关闭，飞书卡片通知",
        "设计：问题状态自动判定（按时关闭/超期关闭/超期未关闭）",
        "流程：问题登记 → 责任人处理 → 问题上升（可选） → 问题关闭"
    ]

    return create_content_slide(
        prs,
        title_cn="重大问题管理模块",
        title_en="Serious Problem Management",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_serious_problem_optimization_slide(prs, src_shapes_xml=None):
    """创建重大问题管理优化页"""
    content_items = [
        "问题：用户连续点击提交按钮导致重复数据入库",
        "方案：自定义@RepeatSubmit注解 + Redis + 拦截器",
        "流程：拦截器判断注解 → 拼接Redis Key → 判断是否重复 → 放行或拒绝",
        "优势：注解式使用，代码复用，支持多模块"
    ]

    return create_content_slide(
        prs,
        title_cn="重大问题管理优化",
        title_en="Serious Problem Optimization",
        section_num="02",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_functional_test_slide(prs, src_shapes_xml=None):
    """创建功能测试页"""
    content_items = [
        "测试范围：四个核心模块的功能测试、边界测试、安全性测试",
        "测试用例：共25个测试用例，覆盖正常流程和异常场景",
        "测试结果：所有用例均符合预期",
        "测试方法：黑盒测试，从实际企业运营场景出发"
    ]

    return create_content_slide(
        prs,
        title_cn="功能测试",
        title_en="Functional Testing",
        section_num="03",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_performance_test_slide(prs, src_shapes_xml=None):
    """创建性能测试页"""
    content_items = [
        "深分页优化：5K页从35736ms降至73.2ms",
        "索引优化：从全表扫描（ALL）优化为范围查询（range）",
        "缓存策略：Redis缓存减少数据库查询压力",
        "测试方法：AOP切面耗时测量，多次执行取平均值"
    ]

    return create_content_slide(
        prs,
        title_cn="性能测试",
        title_en="Performance Testing",
        section_num="03",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_conclusion_slide(prs, src_shapes_xml=None):
    """创建总结页"""
    content_items = [
        "工作成果：完成四个核心模块的设计与实现，系统已在公司内部使用",
        "技术收获：策略模式、延迟关联、索引优化、自定义注解",
        "不足与展望：框架版本升级、功能扩展（部门运营管理、研发效能度量）",
        "未来方向：Spring Boot 3.x + Vue 3 + Java 21虚拟线程"
    ]

    return create_content_slide(
        prs,
        title_cn="总结与展望",
        title_en="Conclusion and Prospects",
        section_num="04",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def create_acknowledgement_slide(prs, src_shapes_xml=None):
    """创建致谢页"""
    content_items = [
        "感谢指导教师的悉心指导",
        "感谢同学和家人的支持",
        "感谢答辩老师的宝贵意见",
        "谢谢大家！"
    ]

    return create_content_slide(
        prs,
        title_cn="致 谢",
        title_en="Acknowledgements",
        section_num="05",
        content_items=content_items,
        src_shapes_xml=src_shapes_xml
    )


def get_slide_shapes_xml(prs, slide_index):
    """获取指定幻灯片的形状XML副本列表

    Args:
        prs: Presentation对象
        slide_index: 幻灯片索引

    Returns:
        形状XML元素的深拷贝列表
    """
    return [deepcopy(shape._element) for shape in prs.slides[slide_index].shapes]


if __name__ == "__main__":
    prs = load_template()

    # Capture template slide shapes before clearing
    cover_shapes = get_slide_shapes_xml(prs, 0)  # 封面
    toc_shapes = get_slide_shapes_xml(prs, 1)    # 目录
    section_shapes = get_slide_shapes_xml(prs, 2) # section标题页

    clear_all_slides(prs)

    # Create all 15 slides
    print("创建封面页...")
    create_cover_slide(prs, cover_shapes)

    print("创建目录页...")
    create_toc_slide(prs, toc_shapes)

    print("创建项目概述页...")
    create_project_overview_slide(prs, section_shapes)

    print("创建系统架构页...")
    create_architecture_slide(prs, section_shapes)

    print("创建BU工时功能页...")
    create_bu_hour_function_slide(prs, section_shapes)

    print("创建BU工时优化页...")
    create_bu_hour_optimization_slide(prs, section_shapes)

    print("创建项目信息管理页...")
    create_project_info_slide(prs, section_shapes)

    print("创建缺陷管理功能页...")
    create_bug_management_function_slide(prs, section_shapes)

    print("创建缺陷管理优化页...")
    create_bug_management_optimization_slide(prs, section_shapes)

    print("创建重大问题管理功能页...")
    create_serious_problem_function_slide(prs, section_shapes)

    print("创建重大问题管理优化页...")
    create_serious_problem_optimization_slide(prs, section_shapes)

    print("创建功能测试页...")
    create_functional_test_slide(prs, section_shapes)

    print("创建性能测试页...")
    create_performance_test_slide(prs, section_shapes)

    print("创建总结页...")
    create_conclusion_slide(prs, section_shapes)

    print("创建致谢页...")
    create_acknowledgement_slide(prs, section_shapes)

    save_ppt(prs, OUTPUT_PATH)
    print(f"\nPPT创建完成！共 {len(prs.slides)} 页")
